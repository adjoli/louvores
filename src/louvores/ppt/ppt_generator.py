from pptx import Presentation

from louvores.db.models import Hino
from louvores.domain.slide_parts import SequenciaHino, TipoParte
from louvores.ppt.layouts import SlideLayout


def _map_layout(tipo: TipoParte) -> SlideLayout:
    match tipo:
        case TipoParte.ESTROFE:
            return SlideLayout.ESTROFE
        case TipoParte.REFRAO:
            return SlideLayout.REFRAO
        case _:
            raise ValueError(f"Tipo de parte desconhecido: {tipo}")


def create_slides(
    hino: Hino, sequencia: SequenciaHino, template_path: str
) -> Presentation:
    prs = Presentation(template_path)

    num_slides = len(sequencia.partes)

    # Slide de título (primeiro slide do template)
    slide_titulo = prs.slides[0]
    slide_titulo.shapes.title.text = hino.titulo
    slide_titulo.shapes[1].text = hino.creditos

    for parte in sequencia.partes:
        layout = _map_layout(parte.tipo)
        slide = prs.slides.add_slide(prs.slide_layouts[layout.value])

        slide.shapes.title.text = hino.titulo  # Título
        slide.shapes[1].text = parte.txt  # Conteúdo

        if len(slide.shapes) > 2:  # Existe espaço para número do slide?
            slide.shapes[2].text = parte.rodape(num_slides)

    return prs
